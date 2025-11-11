import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

def get_image_files_from_folder(folder_path):
    """ 從資料夾取得所有圖片檔案 """
    image_files = []
    for filename in os.listdir(folder_path):
        if filename.endswith(('.png', '.jpg', '.jpeg')):
            image_files.append(filename)
    return image_files

def find_common_images(folders):
    """ 找出三個資料夾中相同檔名的圖片 """
    common_images = set(get_image_files_from_folder(folders[0]))
    for folder in folders[1:]:
        common_images.intersection_update(get_image_files_from_folder(folder))
    return common_images

def add_images_to_slide(presentation, image_paths):
    """ 把圖片加到同一張投影片 """
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # 使用空白投影片
    left = top = Inches(1)  # 圖片的位置
    pic_width = pic_height = Inches(2)  # 圖片的大小
    
    for image_path in image_paths:
        slide.shapes.add_picture(image_path, left, top, pic_width, pic_height)
        left += pic_width + Inches(0.5)  # 調整每個圖片的間隔

def main():
    # 設定三個資料夾路徑
    folder1 = 'C:\\Users\\mike_huang.CLT\\Desktop\\intense_warm\\glow_up_intense_warm_0919'
    folder2 = 'C:\\Users\\mike_huang.CLT\\Desktop\\intense_warm\\glow_up_intense_warm_1016'
    folder3 = 'C:\\Users\\mike_huang.CLT\\Desktop\\intense_warm\\glow_up_intense_warm_1016 - src'

    # 取得所有資料夾中共有的圖片檔名
    common_images = find_common_images([folder1, folder2, folder3])
    
    # 創建一個新的 PowerPoint 演示文稿
    presentation = Presentation()

    # 把每個共有的圖片加到投影片中
    for image_name in common_images:
        # 找出這些圖片的完整路徑
        image_paths = [os.path.join(folder1, image_name),
                       os.path.join(folder2, image_name),
                       os.path.join(folder3, image_name)]
        add_images_to_slide(presentation, image_paths)

    # 儲存 PowerPoint 文件
    presentation.save('output_presentation.pptx')

if __name__ == '__main__':
    main()
